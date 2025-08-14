import json
import requests
import hashlib
import re
from pathlib import Path
from typing import List, Dict, Optional
from urllib.parse import urlparse, unquote
import mimetypes
from dataclasses import dataclass

import PyPDF2
from docx import Document
from bs4 import BeautifulSoup
from pptx import Presentation
import openpyxl

from ..config import settings


@dataclass
class DocumentMetadata:
    source_url: str
    filename: str
    file_type: str
    file_size: int
    download_date: str
    content_hash: str
    title: Optional[str] = None
    description: Optional[str] = None


class DocumentLoader:
    def __init__(self, client_id: str = None):  # ✅ ДОБАВЛЕН client_id
        self.client_id = client_id
        self.session = requests.Session()
        self.session.headers.update({
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
        })

    def load_from_json(self, json_file_path: str) -> List[Dict]:
        """Загружает список документов из JSON файла"""
        with open(json_file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)

        # Поддерживаем разные форматы JSON
        if isinstance(data, list):
            return data
        elif isinstance(data, dict):
            # Ваш формат с result
            if 'result' in data:
                return data['result']
            # Стандартный формат с documents
            elif 'documents' in data:
                return data['documents']
            else:
                raise ValueError("JSON должен содержать 'result' или 'documents'")
        else:
            raise ValueError("Неподдерживаемый формат JSON")

    def download_document(self, url: str, custom_filename: Optional[str] = None) -> Optional[Path]:
        """Скачивает документ по URL"""
        try:
            response = self.session.get(url, timeout=settings.REQUEST_TIMEOUT, stream=True)
            response.raise_for_status()

            # Определяем имя файла
            if custom_filename:
                filename = custom_filename
            else:
                # Пытаемся получить имя из заголовков
                content_disposition = response.headers.get('content-disposition')
                if content_disposition and 'filename=' in content_disposition:
                    filename = content_disposition.split('filename=')[1].strip('"')
                else:
                    # Извлекаем из URL
                    parsed_url = urlparse(url)
                    filename = Path(unquote(parsed_url.path)).name
                    if not filename or '.' not in filename:
                        # Генерируем имя на основе content-type
                        content_type = response.headers.get('content-type', 'application/octet-stream')
                        ext = mimetypes.guess_extension(content_type.split(';')[0]) or '.bin'
                        filename = f"document_{hashlib.md5(url.encode()).hexdigest()[:8]}{ext}"

            # Проверяем расширение
            file_ext = Path(filename).suffix.lower()
            if file_ext not in settings.SUPPORTED_EXTENSIONS:
                print(f"Предупреждение: Неподдерживаемое расширение {file_ext} для {filename}")

            # ✅ ИЗМЕНЕН путь сохранения - используем client_id
            if self.client_id:
                file_path = settings.get_client_temp_documents_dir(self.client_id) / filename
            else:
                # Fallback для совместимости
                temp_dir = settings.DATA_DIR / "temp"
                temp_dir.mkdir(parents=True, exist_ok=True)
                file_path = temp_dir / filename

            # Проверяем размер файла
            content_length = response.headers.get('content-length')
            if content_length and int(content_length) > settings.MAX_FILE_SIZE_MB * 1024 * 1024:
                print(f"Файл {filename} слишком большой ({int(content_length) / 1024 / 1024:.1f}MB)")
                return None

            with open(file_path, 'wb') as f:
                for chunk in response.iter_content(chunk_size=8192):
                    f.write(chunk)

            print(f"Скачан: {filename}")
            return file_path

        except Exception as e:
            print(f"Ошибка при скачивании {url}: {e}")
            return None

    def process_json_documents(self, json_file_path: str) -> List[Dict]:
        """Обрабатывает все документы из JSON файла"""
        documents = self.load_from_json(json_file_path)
        results = []

        for doc_info in documents:
            # ID содержит ссылку на файл
            url = doc_info.get('ID')  # ✅ Берем из ID
            if not url:
                print(f"Пропущен документ без ID: {doc_info}")
                continue

            print(f"Обрабатываем: {url}")

            # Формируем имя файла
            custom_filename = doc_info.get('filename')
            if not custom_filename:
                # Извлекаем имя из URL, убирая лишние символы
                from urllib.parse import urlparse, unquote
                parsed_url = urlparse(url)
                filename = Path(unquote(parsed_url.path)).name
                # Убираем номера версий типа ".1.." из имени
                filename = re.sub(r'\.\d+\.\.', '.', filename)
                custom_filename = filename

            # Скачиваем документ
            file_path = self.download_document(url, custom_filename)

            if file_path and file_path.exists():
                # ИСПРАВЛЕНО: правильно маппим поля
                metadata = {}

                print(f"DEBUG: Исходные данные для {file_path.name}:")
                print(f"       doc_info = {doc_info}")

                # Правильный маппинг полей из 1С JSON
                field_mapping = {
                    'Description': 'description',
                    'Parent': 'parent',
                    'Date': 'date',
                    'GuidDoc': 'guid_doc',
                    'Object_ID': 'object_id',
                    'ID': 'source_url'  # ✅ ID содержит ссылку на файл
                }

                # Сначала применяем правильный маппинг
                for source_field, target_field in field_mapping.items():
                    if source_field in doc_info:
                        metadata[target_field] = doc_info[source_field]
                        print(f"       Маппинг: {source_field} → {target_field} = '{doc_info[source_field]}'")

                print(f"       URL variable: '{url}'")
                print(f"       metadata['source_url']: '{metadata.get('source_url')}'")
                print(f"       doc_info['ID']: '{doc_info.get('ID')}'")

                # Затем добавляем остальные поля в нижнем регистре
                for key, value in doc_info.items():
                    if key not in field_mapping and key not in ['url', 'filename']:
                        normalized_key = key.lower()
                        metadata[normalized_key] = value

                print(f"       ФИНАЛЬНАЯ metadata['source_url']: '{metadata.get('source_url')}'")
                print(f"       ФИНАЛЬНАЯ metadata['description']: '{metadata.get('description')}'")

                # Добавляем стандартные маппинги ТОЛЬКО если их еще нет
                if metadata.get('description') and not metadata.get('title'):
                    metadata['title'] = metadata['description']

                if metadata.get('parent') and not metadata.get('category'):
                    metadata['category'] = metadata['parent']

                print(f"DEBUG: Финальные метаданные для {file_path.name}:")
                print(f"       {len(metadata)} полей: {list(metadata.keys())}")
                print(f"       description = '{metadata.get('description')}'")
                print(f"       title = '{metadata.get('title')}'")
                print(f"       parent = '{metadata.get('parent')}'")
                print(f"       guid_doc = '{metadata.get('guid_doc')}'")

                results.append({
                    'file_path': file_path,
                    'url': url,
                    'metadata': metadata
                })

        return results

class DocumentParser:
    """Парсер для извлечения текста из различных форматов документов"""

    @staticmethod
    def parse_pdf(file_path: Path) -> str:
        """Извлекает текст из PDF"""
        try:
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                return text.strip()
        except Exception as e:
            print(f"Ошибка при парсинге PDF {file_path}: {e}")
            return ""

    @staticmethod
    def parse_docx(file_path: Path) -> str:
        """Извлекает текст из DOCX"""
        try:
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except Exception as e:
            print(f"Ошибка при парсинге DOCX {file_path}: {e}")
            return ""

    @staticmethod
    def parse_txt(file_path: Path) -> str:
        """Читает текстовый файл"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read().strip()
        except UnicodeDecodeError:
            # Пробуем другие кодировки
            for encoding in ['cp1251', 'latin-1']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        return f.read().strip()
                except UnicodeDecodeError:
                    continue
            print(f"Не удалось декодировать файл {file_path}")
            return ""
        except Exception as e:
            print(f"Ошибка при чтении TXT {file_path}: {e}")
            return ""

    @staticmethod
    def parse_html(file_path: Path) -> str:
        """Извлекает текст из HTML"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
                # Убираем скрипты и стили
                for script in soup(["script", "style"]):
                    script.decompose()
                return soup.get_text(separator='\n').strip()
        except Exception as e:
            print(f"Ошибка при парсинге HTML {file_path}: {e}")
            return ""

    @staticmethod
    def parse_pptx(file_path: Path) -> str:
        """Извлекает текст из PPTX"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            print(f"Ошибка при парсинге PPTX {file_path}: {e}")
            return ""

    def parse_document(self, file_path: Path) -> tuple[str, DocumentMetadata]:
        """Парсит документ и возвращает текст с метаданными"""
        file_ext = file_path.suffix.lower()

        # Выбираем парсер
        parsers = {
            '.pdf': self.parse_pdf,
            '.docx': self.parse_docx,
            '.txt': self.parse_txt,
            '.md': self.parse_txt,
            '.html': self.parse_html,
            '.htm': self.parse_html,
            '.pptx': self.parse_pptx,
        }

        parser = parsers.get(file_ext)
        if not parser:
            print(f"Неподдерживаемый формат файла: {file_ext}")
            return "", None

        # Извлекаем текст
        text = parser(file_path)

        # Создаем метаданные
        file_stats = file_path.stat()
        content_hash = hashlib.md5(text.encode()).hexdigest()

        metadata = DocumentMetadata(
            source_url="",  # Будет заполнено позже
            filename=file_path.name,
            file_type=file_ext,
            file_size=file_stats.st_size,
            download_date="",  # Будет заполнено позже
            content_hash=content_hash
        )

        return text, metadata